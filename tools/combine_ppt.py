#! /usr/bin/env python
#
# Combine plot files into a single powerpoint presentation

import geocal
from pptx import Presentation
from pptx.util import Inches

version = "1.0"
usage='''Usage:
  combine_ppt.py [options] <out_name> <in_name>...
  combine_ppt.py -h | --help
  combine_ppt.py -v | --version

This is a short utility to combine a set of plots into one powerpoint 
presentation.

Options:
  -h --help         
       Print this message

  -v --version      
       Print program version
'''

args = geocal.docopt_simple(usage, version=version)
ppt = Presentation()
blank_slide_layout = ppt.slide_layouts[6]
for f in args.in_name:
    slide = ppt.slides.add_slide(blank_slide_layout)
    pic = slide.shapes.add_picture(f, Inches(0), Inches(0), height=Inches(7.5))
    
ppt.save(args.out_name)

